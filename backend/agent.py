import json
import os
import pandas as pd
from langchain_community.chat_message_histories import ChatMessageHistory
from langchain_core.chat_history import BaseChatMessageHistory
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.runnables.history import RunnableWithMessageHistory
from langchain_community.llms import Ollama
from langchain_community.document_loaders import TextLoader  # TextLoader remains the same
import PyPDF2  # To load PDF documents
from docx import Document  # To load .docx files
from pptx import Presentation  # To load .pptx files

class QA_Agent:
    def __init__(self):
        self.llm = Ollama(model="llama3.2")
        print("Model Loaded")

        qa_prompt = ChatPromptTemplate.from_messages([
            ("system", self.get_system_prompt()),
            MessagesPlaceholder("chat_history"),
            ("human", "{input}")
        ])
        q_chain = qa_prompt | self.llm

        self.chat_history = {}
        self.chat_model = RunnableWithMessageHistory(
            q_chain,
            self.get_session_history,
            input_messages_key="input",
            history_messages_key="chat_history"
        )
        self.document_context = []  # Holds the context of documents loaded
        print("Ynjoy!!")

    def get_session_history(self, session_id: str):
        if session_id not in self.chat_history:
            self.chat_history[session_id] = ChatMessageHistory()
        return self.chat_history[session_id]

    def reformat_json(self, json_string, path):
        data = json.loads(json_string)
        formatted_messages = []
        for message in data['messages']:
            if message['type'] == 'human':
                user_content = message['content']
                ai_response = next((msg['content'] for msg in data['messages'] if msg['type'] == 'ai'), None)
                if ai_response:
                    formatted_messages.append({
                        "user_query": user_content,
                        "ai_response": ai_response
                    })

        # Creating the desired JSON structure
        formatted_data = {
            "messages": formatted_messages
        }
        # Convert Python dictionary back to JSON string
        with open(path, 'w') as f:
            json.dump(formatted_data, f, indent=2)
        return formatted_data

    def save_history(self, path="./agent_stup.json"):
        history = self.chat_history["acc_setup"].json()
        self.reformat_json(history, path)
        print(f"History saved to {path}")

    def get_system_prompt(self):
        system_prompt = (
            """
            You are a document-based chatbot. Your task is to help users extract information and insights from uploaded documents. Only answer questions based on the uploaded content. If a question is unrelated or beyond your scope, respond with: "I'm sorry, I can't answer that." Your responses should be concise and accurate. Be polite and helpful at all times
            """
        )
        return system_prompt

    def agent_chat(self, usr_prompt):
        response = self.chat_model.invoke(
            {"input": usr_prompt},
            config={
                "configurable": {"session_id": "acc_setup"}
            }
        )
        return response

    def load_pdf(self, document_path):
        """
        Loads PDF document content using PyPDF2.
        """
        with open(document_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text()
            return text

    def load_csv(self, document_path):
        """
        Loads CSV document content using pandas.
        """
        df = pd.read_csv(document_path)
        text = df.to_string()  # Convert dataframe to a string representation
        return text

    def load_xlsx(self, document_path):
        """
        Loads XLSX document content using pandas.
        """
        df = pd.read_excel(document_path)
        text = df.to_string()  # Convert dataframe to a string representation
        return text

    def load_docx(self, document_path):
        """
        Loads DOCX document content using python-docx.
        """
        doc = Document(document_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text

    def load_cdr_json(self, document_path):
        """
        Loads CDR document content from JSON format.
        """
        try:
            with open(document_path, 'r') as file:
                cdr_data = json.load(file)
            
            # Example: Assuming CDR data is structured with a list of records
            text = ""
            for record in cdr_data.get("calls", []):  # Example structure
                text += f"Call Record: {record.get('caller_id')} to {record.get('receiver_id')}, Duration: {record.get('duration')} mins\n"
            return text
        except Exception as e:
            print(f"Error loading CDR JSON file: {e}")
            return ""

    def load_pptx(self, document_path):
        """
        Loads PPTX document content using python-pptx.
        """
        try:
            presentation = Presentation(document_path)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"Error loading PPTX file: {e}")
            return ""

    def load_document(self, document_path):
        """
        Loads the document content based on its extension.
        Supports text, PDF, CSV, XLSX, DOCX, CDR JSON, and PPTX files.
        """
        if not os.path.exists(document_path):
            print(f"File at {document_path} does not exist!")
            return

        ext = os.path.splitext(document_path)[1].lower()

        if ext == '.txt':
            loader = TextLoader(document_path)
            document = loader.load()
        elif ext == '.pdf':
            document = self.load_pdf(document_path)
            self.document_context = [{"content": document}]
            print(f"Document loaded from {document_path}.")
        elif ext == '.csv':
            document = self.load_csv(document_path)
            self.document_context = [{"content": document}]
            print(f"Document loaded from {document_path}.")
        elif ext == '.xlsx':
            document = self.load_xlsx(document_path)
            self.document_context = [{"content": document}]
            print(f"Document loaded from {document_path}.")
        elif ext == '.docx':
            document = self.load_docx(document_path)
            self.document_context = [{"content": document}]
            print(f"Document loaded from {document_path}.")
        elif ext == '.cdr':
            if document_path.endswith('.json'):
                document = self.load_cdr_json(document_path)
                self.document_context = [{"content": document}]
                print(f"CDR JSON file loaded from {document_path}.")
            else:
                print("Unsupported .cdr file format. Please provide a JSON file.")
                return
        elif ext == '.pptx':
            document = self.load_pptx(document_path)
            self.document_context = [{"content": document}]
            print(f"Document loaded from {document_path}.")
        else:
            print("Unsupported file format. Only text (.txt), PDF (.pdf), CSV (.csv), XLSX (.xlsx), DOCX (.docx), CDR (.cdr), and PPTX (.pptx) are supported.")
            return

    def add_document_to_context(self, document_path):
        self.load_document(document_path)

    def get_document_content(self):
        """
        Retrieves the content of the loaded document.
        """
        if self.document_context:
            return self.document_context[0]["content"]
        return "No document loaded."

def main():
    chat_agent = QA_Agent()
    print("You are now conversing with your assistant. Good Luck!")
    print("You can upload documents by providing the path to the file.")
    
    while True:
        print("Enter your query or the file path to add a document ('/exit' to quit session):")
        prompt = input("Your input: ")

        if prompt == "/exit":
            print("You will have a fortuitous encounter soon, God Speed mate!")
            break

        if os.path.exists(prompt):
            chat_agent.add_document_to_context(prompt)
        else:
            document_content = chat_agent.get_document_content()
            if document_content == "No document loaded.":
                response = chat_agent.agent_chat(prompt)
            else:
                prompt_with_document = f"{document_content}\n\n{prompt}"
                response = chat_agent.agent_chat(prompt_with_document)
            print(f"Assistant: {response}")
            print("-" * 30)

    chat_agent.save_history()

if __name__ == "__main__":
    main()

